from pathlib import Path
from typing import Dict, List

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
SLIDE_CARDS_PATH = PROJECT_ROOT / "docs" / "final_presentation" / "slide_cards.yaml"
OUTPUT_PPTX_PATH = PROJECT_ROOT / "docs" / "final_presentation" / "BidFlow_최종발표_자동생성.pptx"


def _rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _set_textbox_text(text_frame, text: str, font_size: int, bold: bool = False, color: str = "111111") -> None:
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _add_rect(slide, left, top, width, height, fill_hex: str, line_hex: str = "D9D9D9"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(line_hex)
    return shape


def _script_text(script: Dict[str, str]) -> str:
    parts: List[str] = []
    for key in ["opening", "proof", "takeaway", "bridge"]:
        value = script.get(key, "")
        if value:
            parts.append(value)
    return "\n".join(parts)


def _evidence_text(evidence: List[Dict[str, str]]) -> str:
    lines: List[str] = []
    for ev in evidence[:3]:
        metric = ev.get("metric", "")
        value = ev.get("value", "")
        meaning = ev.get("meaning", "")
        lines.append(f"{metric}: {value}")
        if meaning:
            lines.append(f"- {meaning}")
    return "\n".join(lines) if lines else "근거 카드: 추후 입력"


def _chart_map() -> Dict[str, Path]:
    assets = PROJECT_ROOT / "docs" / "final_report_assets"
    return {
        "S06": assets / "chart_01_metric_trends.png",
        "S07": assets / "chart_02_quality_control.png",
        "S08": assets / "chart_03_reliability.png",
        "S09": assets / "chart_04_security_status.png",
    }


def _draw_slide(slide, card: Dict, chart_map: Dict[str, Path]) -> None:
    # Layout constants (16:9)
    # slide size: 13.333 x 7.5
    margin_x = Inches(0.4)
    top_band_h = Inches(0.55)
    content_top = Inches(0.65)
    content_h = Inches(6.45)

    left_col_w = Inches(7.6)
    right_col_x = margin_x + left_col_w + Inches(0.2)
    right_col_w = Inches(4.7)

    # Top one-line conclusion band
    band = _add_rect(slide, margin_x, Inches(0.05), Inches(12.5), top_band_h, "EEF4FF", "C8D9FF")
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_textbox_text(tf, card.get("one_line_conclusion", ""), 16, bold=True, color="1B3F8B")

    # Title
    title_box = slide.shapes.add_textbox(margin_x, content_top, left_col_w, Inches(0.75))
    _set_textbox_text(title_box.text_frame, card.get("title", ""), 30, bold=True)

    # Body copy (2 lines as bullets)
    body_box = slide.shapes.add_textbox(margin_x, content_top + Inches(0.85), left_col_w, Inches(1.6))
    body_tf = body_box.text_frame
    body_tf.clear()
    body_tf.word_wrap = True
    body_lines = card.get("body_copy", [])
    for idx, line in enumerate(body_lines):
        p = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
        p.text = f"• {line}"
        run = p.runs[0]
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(19)
        run.font.color.rgb = _rgb("222222")

    # Plain line (highlight sentence only, no label)
    plain = card.get("plain_line", "")
    plain_shape = _add_rect(
        slide,
        margin_x,
        content_top + Inches(2.55),
        left_col_w,
        Inches(0.95),
        "FFF6E8",
        "F3D5A6",
    )
    plain_tf = plain_shape.text_frame
    plain_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_textbox_text(plain_tf, plain, 18, bold=True, color="7A4A00")

    # Right visual area
    visual_shape = _add_rect(
        slide,
        right_col_x,
        content_top,
        right_col_w,
        Inches(3.95),
        "F9FAFC",
        "D0D5DD",
    )
    visual_tf = visual_shape.text_frame
    visual_tf.vertical_anchor = MSO_ANCHOR.TOP
    _set_textbox_text(visual_tf, card.get("visual", "시각 요소"), 14, bold=False, color="555555")

    slide_id = card.get("id", "")
    chart_path = chart_map.get(slide_id)
    if chart_path and chart_path.exists():
        # Fill most of right visual area with image
        img_left = right_col_x + Inches(0.1)
        img_top = content_top + Inches(0.35)
        img_w = right_col_w - Inches(0.2)
        img_h = Inches(3.45)
        slide.shapes.add_picture(str(chart_path), img_left, img_top, width=img_w, height=img_h)

    # Evidence card
    evidence_shape = _add_rect(
        slide,
        right_col_x,
        content_top + Inches(4.1),
        right_col_w,
        Inches(2.2),
        "F5F7FA",
        "D0D5DD",
    )
    evidence_tf = evidence_shape.text_frame
    evidence_tf.word_wrap = True
    evidence_tf.vertical_anchor = MSO_ANCHOR.TOP
    _set_textbox_text(evidence_tf, _evidence_text(card.get("evidence", [])), 12, color="333333")

    # Footer (id + time)
    footer_box = slide.shapes.add_textbox(margin_x, Inches(7.1), Inches(3.5), Inches(0.25))
    footer_text = f"{card.get('id', '')} | {card.get('time', '')}"
    _set_textbox_text(footer_box.text_frame, footer_text, 11, color="666666")

    # Speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = _script_text(card.get("script", {}))


def build_presentation() -> Path:
    with open(SLIDE_CARDS_PATH, "r", encoding="utf-8") as f:
        payload = yaml.safe_load(f)
    cards: List[Dict] = payload.get("slide_cards", [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    chart_map = _chart_map()
    for card in cards:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _draw_slide(slide, card, chart_map)

    prs.save(str(OUTPUT_PPTX_PATH))
    return OUTPUT_PPTX_PATH


def main() -> None:
    out_path = build_presentation()
    print(f"Generated PPTX: {out_path.as_posix()}")


if __name__ == "__main__":
    main()
